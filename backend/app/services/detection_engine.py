import os
import re
import requests
import urllib.parse
import numpy as np
from bs4 import BeautifulSoup
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity

# ==========================================================
# Lazy Load Model (Prevents Render Boot Crash)
# ==========================================================

_model = None


def get_model():
    global _model
    if _model is None:
        _model = SentenceTransformer(
            "all-MiniLM-L6-v2",
            device="cpu"
        )
    return _model


# ==========================================================
# CENTRAL DETECTION ENGINE
# ==========================================================

class DetectionEngine:
    """
    Complete AI + Plagiarism + Web + File Processor Engine
    """

    def __init__(self):
        self.similarity_threshold = 0.72
        self.max_sentences = 8
        self.max_urls_per_sentence = 3
        self.headers = {
            "User-Agent": "Mozilla/5.0"
        }

    # ======================================================
    # MAIN ENTRY
    # ======================================================

    def analyze_text(self, text):
        if not text or len(text) < 50:
            return self._empty_result()

        sentences = self._split_sentences(text)[:self.max_sentences]

        if not sentences:
            return self._empty_result()

        plagiarism_matches = self._detect_plagiarism(sentences)
        ai_probability = self._detect_ai_generated(text)

        plagiarism_score = (
            max([m["similarity_score"] for m in plagiarism_matches])
            if plagiarism_matches else 0.0
        )

        final_score = round(
            (plagiarism_score * 0.7) + (ai_probability * 0.3),
            2
        )

        return {
            "plagiarism_matches": plagiarism_matches,
            "plagiarism_score": plagiarism_score,
            "ai_probability": ai_probability,
            "final_score": final_score
        }

    # ======================================================
    # PLAGIARISM DETECTION
    # ======================================================

    def _detect_plagiarism(self, sentences):
        model = get_model()
        sentence_embeddings = model.encode(sentences)

        matches = []

        for idx, sentence in enumerate(sentences):

            urls = self._search_text(sentence)

            for url in urls:
                content = self._fetch_content(url)

                if not content:
                    continue

                content = content[:5000]

                try:
                    content_embedding = model.encode([content])[0]

                    similarity = cosine_similarity(
                        [sentence_embeddings[idx]],
                        [content_embedding]
                    )[0][0]

                    if similarity >= self.similarity_threshold:
                        matches.append({
                            "source_url": url,
                            "matched_text": sentence,
                            "similarity_score": float(similarity)
                        })

                except Exception:
                    continue

        return self._deduplicate_matches(matches)

    # ======================================================
    # AI DETECTION (Pattern Based)
    # ======================================================

    def _detect_ai_generated(self, text):
        sentences = self._split_sentences(text)

        if len(sentences) < 3:
            return 0.0

        sentence_lengths = [len(s.split()) for s in sentences]
        variance = np.var(sentence_lengths)

        if variance < 5:
            return 0.75
        elif variance < 10:
            return 0.55
        else:
            return 0.25

    # ======================================================
    # WEB SEARCH (DuckDuckGo)
    # ======================================================

    def _search_text(self, query):
        results = []

        try:
            encoded_query = urllib.parse.quote(query[:200])
            url = f"https://html.duckduckgo.com/html/?q={encoded_query}"

            response = requests.get(
                url,
                headers=self.headers,
                timeout=10
            )
            response.raise_for_status()

            soup = BeautifulSoup(response.text, "html.parser")
            links = soup.find_all("a", class_="result__a")

            for link in links[:self.max_urls_per_sentence]:
                href = link.get("href")
                if href:
                    results.append(href)

        except Exception:
            pass

        return results

    # ======================================================
    # FETCH WEB CONTENT
    # ======================================================

    def _fetch_content(self, url):
        try:
            response = requests.get(
                url,
                headers=self.headers,
                timeout=10
            )
            response.raise_for_status()

            soup = BeautifulSoup(response.content, "html.parser")

            for script in soup(["script", "style"]):
                script.decompose()

            text = soup.get_text(separator=" ")
            text = " ".join(text.split())

            return text[:15000]

        except Exception:
            return None

    # ======================================================
    # FILE PROCESSOR
    # ======================================================

    def extract_text_from_file(self, file_path, file_type):

        if file_type == "txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        if file_type == "pdf":
            import pdfplumber
            text = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text.append(t)
            return "\n".join(text)

        if file_type in ["docx", "doc"]:
            from docx import Document
            doc = Document(file_path)
            return "\n".join(
                para.text for para in doc.paragraphs
                if para.text.strip()
            )

        if file_type in ["pptx", "ppt"]:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return "\n".join(text)

        raise ValueError("Unsupported file type")

    # ======================================================
    # HELPERS
    # ======================================================

    def _split_sentences(self, text):
        sentences = re.split(r"[.!?]+", text)
        return [s.strip() for s in sentences if len(s.strip()) > 20]

    def _deduplicate_matches(self, matches):
        if not matches:
            return []

        matches = sorted(
            matches,
            key=lambda m: m["similarity_score"],
            reverse=True
        )

        unique = []
        for match in matches:
            if match not in unique:
                unique.append(match)

        return unique[:25]

    def _empty_result(self):
        return {
            "plagiarism_matches": [],
            "plagiarism_score": 0.0,
            "ai_probability": 0.0,
            "final_score": 0.0
        }
