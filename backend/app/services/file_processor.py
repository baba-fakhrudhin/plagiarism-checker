import os

class FileProcessor:
    """Handle file extraction from various formats"""
    
    def extract_text(self, file_path, file_type):
        """Extract text from various file formats"""
        try:
            if file_type == 'txt':
                return self._extract_txt(file_path)
            elif file_type == 'pdf':
                return self._extract_pdf(file_path)
            elif file_type in ['docx', 'doc']:
                return self._extract_docx(file_path)
            elif file_type in ['pptx', 'ppt']:
                return self._extract_pptx(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            raise Exception(f"Error extracting text from {file_type}: {str(e)}")
    
    def _extract_txt(self, file_path):
        """Extract text from TXT files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            raise Exception(f"Error extracting TXT: {str(e)}")
    
    def _extract_pdf(self, file_path):
        """Extract text from PDF files"""
        try:
            import pdfplumber
            text = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
            return '\n'.join(text) if text else ""
        except Exception as e:
            raise Exception(f"Error extracting PDF: {str(e)}")
    
    def _extract_docx(self, file_path):
        """Extract text from DOCX/DOC files"""
        try:
            from docx import Document
            doc = Document(file_path)
            text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text.append(para.text)
            return '\n'.join(text) if text else ""
        except Exception as e:
            raise Exception(f"Error extracting DOCX: {str(e)}")
    
    def _extract_pptx(self, file_path):
        """Extract text from PPTX/PPT files"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return '\n'.join(text) if text else ""
        except Exception as e:
            raise Exception(f"Error extracting PPTX: {str(e)}")